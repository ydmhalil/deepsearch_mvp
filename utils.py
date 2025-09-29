import os
from typing import Optional


def iter_files(root_dir: str, exts=None):
    exts = exts or ['.pdf', '.docx', '.txt', '.pptx', '.xlsx']
    for dirpath, dirs, files in os.walk(root_dir):
        for f in files:
            if any(f.lower().endswith(e) for e in exts):
                yield os.path.join(dirpath, f)

def extract_text_from_pdf(path: str) -> str:
    text_parts = []
    try:
        import pdfplumber
    except Exception:
        # pdfplumber not installed
        return ''
    try:
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                # Normal metin çıkarma
                text = page.extract_text() or ''
                text_parts.append(text)
                
                # Eğer sayfa çok az metin içeriyorsa OCR dene
                if len(text.strip()) < 50:
                    try:
                        # OCR ile metin çıkarma
                        ocr_text = extract_text_with_ocr_from_page(page)
                        if ocr_text.strip():
                            text_parts.append(f"[OCR] {ocr_text}")
                    except Exception:
                        pass  # OCR başarısız olursa sessizce devam et
                        
    except Exception:
        return ''
    return '\n'.join(text_parts)


def extract_text_with_ocr_from_page(page) -> str:
    """PDF sayfasından OCR ile metin çıkarma"""
    try:
        import pytesseract
        from PIL import Image
        import io
        
        # Tesseract path kontrolü - Windows'ta genelde şuralarda:
        possible_paths = [
            r'C:\Program Files\Tesseract-OCR\tesseract.exe',
            r'C:\Program Files (x86)\Tesseract-OCR\tesseract.exe',
            'tesseract'  # PATH'de varsa
        ]
        
        tesseract_cmd = None
        for path in possible_paths:
            try:
                import subprocess
                subprocess.run([path, '--version'], capture_output=True, check=True)
                tesseract_cmd = path
                break
            except (subprocess.CalledProcessError, FileNotFoundError):
                continue
        
        if tesseract_cmd:
            pytesseract.pytesseract.tesseract_cmd = tesseract_cmd
        else:
            # Tesseract bulunamadıysa OCR'ı atla
            return ''
            
    except Exception:
        return ''
    
    try:
        # PDF sayfasını görüntüye çevir
        im = page.within_bbox(page.bbox).to_image()
        # PIL Image'a çevir
        pil_image = im.original
        # OCR uygula (Türkçe dil desteği)
        ocr_text = pytesseract.image_to_string(pil_image, lang='tur+eng')
        return ocr_text
    except Exception:
        return ''


def extract_text_from_docx(path: str) -> str:
    try:
        from docx import Document
    except Exception:
        return ''
    try:
        doc = Document(path)
        return '\n'.join(p.text for p in doc.paragraphs)
    except Exception:
        return ''


def extract_text_from_pptx(path: str) -> str:
    """PowerPoint dosyalarından metin çıkarma"""
    try:
        from pptx import Presentation
    except Exception:
        return ''
    try:
        prs = Presentation(path)
        text_parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_parts.append(shape.text)
        return '\n'.join(text_parts)
    except Exception:
        return ''


def extract_text_from_xlsx(path: str) -> str:
    """Excel dosyalarından metin çıkarma"""
    try:
        from openpyxl import load_workbook
    except Exception:
        return ''
    try:
        workbook = load_workbook(path, data_only=True)
        text_parts = []
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            # Sheet adını ekle
            text_parts.append(f"Sheet: {sheet_name}")
            for row in sheet.iter_rows(values_only=True):
                row_text = []
                for cell in row:
                    if cell is not None and str(cell).strip():
                        row_text.append(str(cell))
                if row_text:
                    text_parts.append(' | '.join(row_text))
        return '\n'.join(text_parts)
    except Exception:
        return ''


def is_binary_file(path: str) -> bool:
    try:
        with open(path, 'rb') as f:
            chunk = f.read(1024)
            if b'\0' in chunk:
                return True
    except Exception:
        return False
    return False


def extract_text(path: str) -> str:
    ext = os.path.splitext(path)[1].lower()
    if ext == '.pdf':
        return extract_text_from_pdf(path)
    elif ext == '.docx':
        return extract_text_from_docx(path)
    elif ext == '.pptx':
        return extract_text_from_pptx(path)
    elif ext == '.xlsx':
        return extract_text_from_xlsx(path)
    elif ext == '.txt':
        try:
            # skip binary files saved as .txt
            if is_binary_file(path):
                return ''
            with open(path, 'r', encoding='utf-8') as f:
                return f.read()
        except Exception:
            try:
                with open(path, 'r', encoding='latin-1') as f:
                    return f.read()
            except Exception:
                return ''
    else:
        return ''
